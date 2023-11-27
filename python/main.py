import customtkinter
from tkinter import filedialog
from pptx import Presentation
from collections import Counter

class PowerPointProcessor:
    def __init__(self, file_path, textbox, button_summurai):
        self.file_path = file_path
        self.textbox = textbox
        self.button_summurai = button_summurai

    def set_appearance_and_theme(self):
        customtkinter.set_appearance_mode("dark")
        customtkinter.set_default_color_theme("green")

    def get_font_size(self, paragraph):
        font_size = None
        if paragraph.runs and paragraph.runs[0].font.size:
            font_size = paragraph.runs[0].font.size.pt
        return font_size

    def most_common_font_size(self, presentation):
        font_sizes = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        font_size = self.get_font_size(paragraph)
                        if font_size is not None:
                            font_sizes.append(font_size)

        if font_sizes:
            most_common_size = Counter(font_sizes).most_common(1)[0][0]
            return most_common_size
        else:
            return None

    def simple_summarization(self, text):
        # Basic summarization: Take the first 50 words
        words = text.split()
        summary = ' '.join(words[:50])
        return summary

    def process_powerpoint_file(self):
        formatted_lines = []
        output_text = ""
        prs = Presentation(self.file_path)

        # Calculate the most common font size
        common_font_size = self.most_common_font_size(prs)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text
                        font_size = self.get_font_size(paragraph)
                        line_info = {
                            "text": text,
                            "is_key_text": False,
                            "is_bold": run.font.bold,
                            "is_italic": run.font.italic,
                            "is_underline": run.font.underline,
                        }

                        if font_size is not None and font_size > 1.1 * common_font_size:
                            line_info["is_key_text"] = True

                        if shape == slide.shapes[0]:
                            line_info["is_key_text"] = True

                        formatted_lines.append(line_info)

        for line_info in formatted_lines:
            text = line_info["text"]
            is_key_text = line_info["is_key_text"]
            is_bold = line_info["is_bold"]
            is_italic = line_info["is_italic"]
            is_underline = line_info["is_underline"]

            output_text += "\n"
            if is_key_text:
                output_text += f"Key Text: {text}\n"
            else:
                output_text += f"Text: {text}\n"

            if is_bold:
                output_text += "  - Bold\n"
            if is_italic:
                output_text += "  - Italic\n"
            if is_underline:
                output_text += "  - Underline\n"

        # Perform simple summarization
        summary = self.simple_summarization(output_text)
        
        # Display the summary in the textbox
        self.textbox.delete("1.0", "end")
        self.textbox.insert("1.0", summary)

        # Print a message about the summarization
        print("Textfile successfully summuraised!")

        if self.button_summurai[0]:
            self.button_summurai[0].configure(state="normal", command=lambda: self.process_powerpoint_file())

class GUI:
    def __init__(self):
        self.set_appearance_and_theme()

        self.root = customtkinter.CTk()
        self.root.geometry("1200x800")

        self.frame = customtkinter.CTkFrame(master=self.root)
        self.frame.pack(pady=0, padx=0, fill="both", expand=True)

        self.label = customtkinter.CTkLabel(master=self.frame, text="Please select a .pptx file", font=("Helvetica", 24))
        self.label.pack(pady=20, padx=0)

        self.button_summurai = [None]  
        self.button = customtkinter.CTkButton(master=self.frame, text="Select PowerPoint File", command=self.select_pptx_file)
        self.button.pack(pady=20, padx=0)
        print("Select PowerPoint File button created")

        self.textbox = customtkinter.CTkTextbox(master=self.frame, font=("Helvetica", 14), height=400, width=800)
        self.textbox.pack(pady=10, padx=0)

        self.root.mainloop()

    def set_appearance_and_theme(self):
        customtkinter.set_appearance_mode("dark")
        customtkinter.set_default_color_theme("green")

    def select_pptx_file(self):
        file_path = filedialog.askopenfilename(
            title="Select a .pptx file",
            filetypes=[("PowerPoint files", "*.pptx")]
        )
        if file_path:
            self.label.configure(text='SUCCESS: Selected PowerPoint file', font=("Helvetica", 24))
            self.create_summurai_button(file_path)
            if self.button_summurai[0]:
                self.button_summurai[0].configure(state="normal", command=lambda: self.process_powerpoint_file(file_path))
        else:
            self.label.configure(text='No file selected', font=("Helvetica", 24))

    def create_summurai_button(self, file_path):
        if not self.button_summurai[0]:
            button = customtkinter.CTkButton(master=self.frame, text="SummurAI", state="disabled")
            button.pack(pady=20, padx=0)
            self.button_summurai[0] = button
            print("SummurAI button created")
        else:
            self.button_summurai[0].configure(command=lambda: self.process_powerpoint_file(file_path))

    def process_powerpoint_file(self, file_path):
        processor = PowerPointProcessor(file_path, self.textbox, self.button_summurai)
        processor.process_powerpoint_file()

if __name__ == "__main__":
    gui = GUI()



