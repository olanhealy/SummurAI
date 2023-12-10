from pptx import Presentation
from collections import Counter

# Helper functions ---------------------------------------------------

# Get font size of a paragraph
def get_font_size(paragraph):
    font_size = None
    if paragraph.runs and paragraph.runs[0].font.size:
        font_size = paragraph.runs[0].font.size.pt
    return font_size

# Get the most common font size in a PowerPoint presentation
def most_common_font_size(presentation):
    font_sizes = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    font_size = get_font_size(paragraph)
                    if font_size is not None:
                        font_sizes.append(font_size)

    if font_sizes:
        most_common_size = Counter(font_sizes).most_common(1)[0][0]
        return most_common_size
    else:
        return None  # Return None if there are no font sizes
    
# Main function ------------------------------------------------------

# Process the PowerPoint presentation and return the formatted information
def process_powerpoint(prs):
    # Calculate the most common font size
    common_font_size = most_common_font_size(prs)

    formatted_lines = []

    # Iterate through the slides taking the info such as font size/bold/italic/underline
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    font_size = get_font_size(paragraph)
                    line_info = {
                        "text": text,
                        "is_key_text": False,
                        "is_bold": run.font.bold,
                        "is_italic": run.font.italic,
                        "is_underline": run.font.underline,
                    }

                    # If the font size is larger than 1.1 times the most common font size,
                    if font_size is not None and font_size > 1.1 * common_font_size:
                        line_info["is_key_text"] = True

                    # If the shape is the first shape in the slide,
                    if shape == slide.shapes[0]:
                        line_info["is_key_text"] = True

                    formatted_lines.append(line_info)
    
    return formatted_lines

