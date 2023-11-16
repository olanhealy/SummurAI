import os
import customtkinter
from tkinter import filedialog
from pptx import Presentation  # Import the python-pptx library

def set_appearance_and_theme():
    customtkinter.set_appearance_mode("dark")
    customtkinter.set_default_color_theme("green")

def process_powerpoint_file(file_path):
    formatted_lines = []
    default_font_size = 12

    prs = Presentation(file_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    line_info = {
                        "text": text,
                        "is_heading": False,
                        "is_bold": run.font.bold,
                        "is_italic": run.font.italic,
                        "is_underline": run.font.underline,
                    }

                    font_size = run.font.size.pt if run.font.size is not None else default_font_size
                    if font_size > 18:
                        line_info["is_heading"] = True
                    formatted_lines.append(line_info)

    output_file_path = "formatted_lines.txt"

    with open(output_file_path, "w", encoding="utf-8") as output_file:
        for line_info in formatted_lines:
            text = line_info["text"]
            is_heading = line_info["is_heading"]
            is_bold = line_info["is_bold"]
            is_italic = line_info["is_italic"]
            is_underline = line_info["is_underline"]

            output_file.write("\n")
            if is_heading:
                output_file.write(f"H: {text}\n")
            else:
                output_file.write(f"T: {text}\n")

            if is_bold:
                output_file.write("  - B\n")
            if is_italic:
                output_file.write("  - I\n")
            if is_underline:
                output_file.write("  - U\n")

    print(f"PowerPoint file processed. Output saved to: {output_file_path}")

def select_pptx_file():
    file_path = filedialog.askopenfilename(
        title="Select a .pptx file",
        filetypes=[("PowerPoint files", "*.pptx")]
    )
    if file_path:
        label.configure(text='SUCCESS: Selected PowerPoint file', font=("Helvetica", 24))
        process_powerpoint_file(file_path)
        create_summurai_button()
    else:
        label.configure(text='No file selected', font=("Helvetica", 24))

def create_summurai_button():
    button_summurai = customtkinter.CTkButton(master=frame, text="SummurAI", command=placeholder)
    button_summurai.pack(pady=20, padx=0)
    print("SummurAI button created")

def placeholder():
    print("Debug")
    # Add your code here

def main():
    set_appearance_and_theme()

    global root
    root = customtkinter.CTk()
    root.geometry("1200x800")

    global frame
    frame = customtkinter.CTkFrame(master=root)
    frame.pack(pady=0, padx=0, fill="both", expand=True)

    global label
    label = customtkinter.CTkLabel(master=frame, text="Please select a .pptx file", font=("Helvetica", 24))
    label.pack(pady=20, padx=0)

    global button
    button = customtkinter.CTkButton(master=frame, text="Select PowerPoint File", command=select_pptx_file)
    button.pack(pady=20, padx=0)
    print("Select PowerPoint File button created")

    global textbox
    textbox = customtkinter.CTkTextbox(master=frame, font=("Helvetica", 14), height=800, width=800)
    textbox.pack(pady=10, padx=0)

    root.mainloop()

if __name__ == "__main__":
    main()
