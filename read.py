from pptx import Presentation
from collections import Counter

def get_font_size(paragraph):
    """
    Get the font size of a paragraph in points.
    """
    font_size = None
    if paragraph.runs and paragraph.runs[0].font.size:
        font_size = paragraph.runs[0].font.size.pt
    return font_size

def most_common_font_size(presentation):
    """
    Find the most common font size in a PowerPoint presentation.
    """
    font_sizes = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    font_size = get_font_size(paragraph)
                    if font_size is not None:
                        font_sizes.append(font_size)

    if font_sizes:
        most_common_size = Counter(font_sizes).most_common(1)[0][0]
        return most_common_size
    else:
        return None  # Return None if there are no font sizes

prs = Presentation("pptx_files/CI_Pipeline.pptx")

# Define a list to store structured formatting information for each line
formatted_lines = []

# Specify the output file path
output_file_path = "formatted_lines.txt"

# Calculate the most common font size
common_font_size = most_common_font_size(prs)

with open(output_file_path, "w", encoding="utf-8") as output_file:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    font_size = get_font_size(paragraph)
                    line_info = {
                        "text": text,
                        "is_key_text": False,
                        "is_bold": run.font.bold,
                        "is_italic": run.font.italic,
                        "is_underline": run.font.underline,
                    }

                    # Check if the line is Key Text based on font size
                    if font_size is not None and font_size > 1.1 * common_font_size:
                        line_info["is_key_text"] = True

                    # Check if the line is a title based on its position in the slide
                    if shape == slide.shapes[0]:
                        line_info["is_key_text"] = True

                    formatted_lines.append(line_info)

    for line_info in formatted_lines:
        text = line_info["text"]
        is_key_text = line_info["is_key_text"]
        is_bold = line_info["is_bold"]
        is_italic = line_info["is_italic"]
        is_underline = line_info["is_underline"]

        # Write the information to the output file
        with open(output_file_path, "a", encoding="utf-8") as output_file:
            output_file.write("\n")
            if is_key_text:
                output_file.write(f"Key Text: {text}\n")
            else:
                output_file.write(f"Text: {text}\n")

            if is_bold:
                output_file.write("  - Bold\n")
            if is_italic:
                output_file.write("  - Italic\n")
            if is_underline:
                output_file.write("  - Underline\n")

print(f"Formatted lines information has been saved to {output_file_path}")
